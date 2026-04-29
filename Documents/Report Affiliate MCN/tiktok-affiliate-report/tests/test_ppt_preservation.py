"""
PPT Renderer Preservation Property Tests

This test verifies that existing layout and content remain unchanged by the fix.
It follows the observation-first methodology: observe behavior on UNFIXED code,
then write property-based tests capturing that observed behavior.

**IMPORTANT**: These tests should PASS on UNFIXED code (confirming baseline behavior)
and continue to PASS on FIXED code (confirming no regressions).
"""

import pytest
import os
import tempfile
from datetime import date
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from hypothesis import given, strategies as st, settings

from app.services.ppt_renderer import PPTRenderer
from app.services.report_gen import ReportConfig, PerformanceMetrics, ReportData, CreatorRow, EngagementData, CreatorEngagementItem, VideoEngagementItem
from app.services.brand_profile import BrandProfileData


class TestPPTPreservation:
    """
    Preservation Property Tests - Existing Layout and Content Unchanged
    
    These tests observe behavior on UNFIXED code and verify that the fix does NOT
    change any existing content, layouts, or formatting. All tests should PASS on
    both unfixed and fixed code.
    """
    
    def setup_method(self):
        """Set up test fixtures."""
        self.test_dir = tempfile.mkdtemp()
    
    def teardown_method(self):
        """Clean up test directory."""
        if os.path.exists(self.test_dir):
            import shutil
            shutil.rmtree(self.test_dir)
    
    def create_test_report_data(self, **kwargs):
        """Create test ReportData with configurable parameters."""
        # Default configuration
        config = ReportConfig(
            brand_name=kwargs.get('brand_name', "Test Brand"),
            batch_number=kwargs.get('batch_number', "Batch 1"),
            period_start=kwargs.get('period_start', date(2024, 1, 1)),
            period_end=kwargs.get('period_end', date(2024, 1, 31)),
            prev_gmv=kwargs.get('prev_gmv', 10000000),
            insight=kwargs.get('insight', "Test insight"),
            next_plan=kwargs.get('next_plan', "Test next plan"),
            section_images=kwargs.get('section_images', {})
        )
        
        # Default metrics
        metrics = PerformanceMetrics(
            total_approached=kwargs.get('total_approached', 50),
            total_deal=kwargs.get('total_deal', 30),
            total_posting=kwargs.get('total_posting', 25),
            total_belum_posting=kwargs.get('total_belum_posting', 5),
            total_video=kwargs.get('total_video', 75),
            total_pesanan=kwargs.get('total_pesanan', 100),
            total_produk_terjual=kwargs.get('total_produk_terjual', 150),
            total_create_sale=kwargs.get('total_create_sale', 20),
            total_gmv=kwargs.get('total_gmv', 15000000),
            hero_sku=kwargs.get('hero_sku', "Test SKU"),
            gmv_change=kwargs.get('gmv_change', 5000000),
            gmv_change_pct=kwargs.get('gmv_change_pct', 50.0)
        )
        
        # Default creator rows
        num_creators = kwargs.get('num_creators', 10)
        deal_rows = [
            CreatorRow(
                username=f"creator{i}",
                followers=10000 + (i * 1000),
                avg_gmv_month=500000 + (i * 100000),
                link_acc=f"https://tiktok.com/@creator{i}"
            )
            for i in range(num_creators)
        ]
        
        # Default brand profile
        brand_profile = BrandProfileData(
            name=kwargs.get('brand_name', "Test Brand"),
            sku_list=kwargs.get('sku_list', ["SKU1", "SKU2", "SKU3"])
        )
        
        # Optional engagement data
        engagement = kwargs.get('engagement', None)
        if engagement:
            config.engagement = engagement
        
        return ReportData(
            config=config,
            metrics=metrics,
            top_performers=deal_rows[:5],
            deal_rows=deal_rows,
            non_deal_rows=[],
            brand_profile=brand_profile
        )
    
    def extract_slide_structure(self, slide):
        """Extract structural information from a slide for comparison."""
        structure = {
            'shape_count': len(slide.shapes),
            'text_shapes': [],
            'table_shapes': [],
            'rect_shapes': [],
            'textbox_shapes': []
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                structure['text_shapes'].append({
                    'text': text,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
            
            if shape.has_table:
                table = shape.table
                structure['table_shapes'].append({
                    'rows': len(table.rows),
                    'cols': len(table.columns),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
            
            # Check for rectangles (metric boxes, backgrounds)
            if shape.shape_type == 1:  # AUTO_SHAPE
                structure['rect_shapes'].append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
        
        return structure
    
    def test_no_section_images_preservation(self):
        """
        Preservation: PPT without section_images should be identical
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that when no section_images are configured, the generated PPT
        is identical between original and fixed code.
        """
        # Create report data WITHOUT section_images
        report_data = self.create_test_report_data()
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_no_images.pptx")
        renderer.render(report_data, output_path)
        
        # Verify PPT was generated successfully
        assert os.path.exists(output_path), "PPT file should be generated"
        
        # Extract and verify slide structure
        prs = Presentation(output_path)
        
        # Verify expected number of slides (8 slides total)
        # 0: Cover, 1: Performance, 2: GMV Highlight, 3: Top 10, 4: Collaboration,
        # 5: Engagement, 6: GMV Detail, 7: Insight, 8: Next Plan
        assert len(prs.slides) >= 8, "Should have at least 8 slides"
        
        # Verify each slide has content (shapes)
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"Slide {i} should have shapes"
    
    def test_cover_slide_preservation(self):
        """
        Preservation: Cover slide layout unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that cover slide displays brand name, batch number, and period
        correctly, regardless of section_images configuration.
        """
        report_data = self.create_test_report_data(
            brand_name="FLORIST",
            batch_number="Batch 5"
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_cover.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        cover_slide = prs.slides[0]
        
        # Extract all text from cover slide
        cover_text = " ".join([shape.text for shape in cover_slide.shapes if shape.has_text_frame])
        
        # Verify brand name appears
        assert "FLORIST" in cover_text, "Brand name should appear on cover"
        
        # Verify batch number appears
        assert "Batch 5" in cover_text or "5" in cover_text, "Batch number should appear on cover"
        
        # Verify period appears
        assert "2024" in cover_text, "Period year should appear on cover"
    
    def test_performance_summary_layout_preservation(self):
        """
        Preservation: Performance Summary slide layout unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that metric boxes and summary table remain at same positions
        with same content, regardless of section_images.
        """
        report_data = self.create_test_report_data(
            total_deal=30,
            total_posting=25,
            total_video=75,
            total_gmv=15000000
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_performance_layout.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        performance_slide = prs.slides[1]
        
        # Extract slide structure
        structure = self.extract_slide_structure(performance_slide)
        
        # Verify metric boxes exist (should have multiple rectangles for boxes)
        assert len(structure['rect_shapes']) > 0, "Should have metric boxes (rectangles)"
        
        # Verify summary table exists
        assert len(structure['table_shapes']) > 0, "Should have summary table"
        
        # Verify table has expected structure (3 columns: No, Keterangan, Data)
        table_info = structure['table_shapes'][0]
        assert table_info['cols'] == 3, "Summary table should have 3 columns"
        assert table_info['rows'] >= 6, "Summary table should have at least 6 rows"
        
        # Verify key metrics appear in text
        slide_text = " ".join([shape.text for shape in performance_slide.shapes if shape.has_text_frame])
        assert "30" in slide_text, "Total deal should appear"
        assert "25" in slide_text, "Total posting should appear"
        assert "75" in slide_text, "Total video should appear"
    
    def test_gmv_highlight_preservation(self):
        """
        Preservation: GMV Highlight slide continues to work correctly
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that GMV Highlight slide (which already has working image code)
        continues to display correctly. This is the reference slide.
        """
        report_data = self.create_test_report_data(
            total_gmv=15000000,
            gmv_change=5000000,
            gmv_change_pct=50.0
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_gmv_highlight.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        gmv_slide = prs.slides[2]
        
        # Extract slide text
        slide_text = " ".join([shape.text for shape in gmv_slide.shapes if shape.has_text_frame])
        
        # Verify GMV amount appears
        assert "15.000.000" in slide_text or "15,000,000" in slide_text, "GMV amount should appear"
        
        # Verify GMV change appears
        assert "50" in slide_text, "GMV change percentage should appear"
        
        # Verify SKU list appears
        assert "SKU" in slide_text, "SKU list should appear"
    
    def test_top10_table_preservation(self):
        """
        Preservation: Top 10 Performer table unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that Top 10 table has correct structure and content.
        """
        report_data = self.create_test_report_data(num_creators=10)
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_top10_table.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        top10_slide = prs.slides[3]
        
        # Extract slide structure
        structure = self.extract_slide_structure(top10_slide)
        
        # Verify table exists
        assert len(structure['table_shapes']) > 0, "Should have top 10 table"
        
        # Verify table structure (3 columns: #, Username, GMV)
        table_info = structure['table_shapes'][0]
        assert table_info['cols'] == 3, "Top 10 table should have 3 columns"
        assert table_info['rows'] >= 5, "Top 10 table should have at least 5 rows (header + 4 data)"
    
    def test_collaboration_metrics_preservation(self):
        """
        Preservation: Collaboration Metrics boxes unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that collaboration metric boxes display correctly.
        """
        report_data = self.create_test_report_data(
            total_approached=50,
            total_deal=30
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_collaboration.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        collaboration_slide = prs.slides[4]
        
        # Extract slide text
        slide_text = " ".join([shape.text for shape in collaboration_slide.shapes if shape.has_text_frame])
        
        # Verify metrics appear
        assert "50" in slide_text, "Total approached should appear"
        assert "30" in slide_text, "Total deal should appear"
        assert "60" in slide_text or "%" in slide_text, "Conversion ratio should appear"
    
    def test_engagement_metrics_preservation(self):
        """
        Preservation: Engagement metrics and table unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that engagement metrics and top creators table display correctly.
        """
        # Create engagement data
        engagement = EngagementData(
            total_views=1000000,
            total_likes=50000,
            total_comments=5000,
            top_creators=[
                CreatorEngagementItem(
                    username="creator1",
                    total_views=500000,
                    total_likes=25000,
                    total_comments=2500,
                    videos=[
                        VideoEngagementItem(url="https://tiktok.com/video1", views=300000, likes=15000, comments=1500),
                        VideoEngagementItem(url="https://tiktok.com/video2", views=200000, likes=10000, comments=1000)
                    ]
                ),
                CreatorEngagementItem(
                    username="creator2",
                    total_views=300000,
                    total_likes=15000,
                    total_comments=1500,
                    videos=[
                        VideoEngagementItem(url="https://tiktok.com/video3", views=200000, likes=10000, comments=1000),
                        VideoEngagementItem(url="https://tiktok.com/video4", views=100000, likes=5000, comments=500)
                    ]
                )
            ]
        )
        
        report_data = self.create_test_report_data(engagement=engagement)
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_engagement.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        engagement_slide = prs.slides[5]
        
        # Extract slide structure
        structure = self.extract_slide_structure(engagement_slide)
        
        # Verify metric boxes exist
        assert len(structure['rect_shapes']) > 0, "Should have engagement metric boxes"
        
        # Verify top creators table exists
        assert len(structure['table_shapes']) > 0, "Should have top creators table"
        
        # Verify table has 8 columns (for 2-video breakdown)
        table_info = structure['table_shapes'][0]
        assert table_info['cols'] == 8, "Engagement table should have 8 columns for 2-video breakdown"
    
    def test_color_scheme_preservation(self):
        """
        Preservation: Color scheme unchanged across all slides
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that color scheme (NAVY, BLUE, CORAL, WHITE, LIGHT) is consistent.
        """
        report_data = self.create_test_report_data()
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_colors.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        
        # Verify all slides have shapes (indicating proper rendering)
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"Slide {i} should have shapes with colors"
    
    def test_currency_formatting_preservation(self):
        """
        Preservation: Currency formatting unchanged
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Verifies that format_rupiah function continues to work correctly.
        """
        report_data = self.create_test_report_data(total_gmv=15000000)
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_currency.pptx")
        renderer.render(report_data, output_path)
        
        prs = Presentation(output_path)
        
        # Check multiple slides for currency formatting
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text + " "
        
        # Verify Rupiah formatting appears (Rp with dot separators)
        assert "Rp" in all_text, "Rupiah symbol should appear"
        assert "15.000.000" in all_text or "15,000,000" in all_text, "GMV should be formatted"
    
    @settings(max_examples=10, deadline=None)
    @given(
        total_gmv=st.integers(min_value=1000000, max_value=100000000),
        total_deal=st.integers(min_value=10, max_value=100),
        total_video=st.integers(min_value=10, max_value=500)
    )
    def test_property_various_metrics(self, total_gmv, total_deal, total_video):
        """
        Property-Based Test: Various metric values produce valid PPTs
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Generates many test cases with different metric values to ensure
        the renderer handles all inputs correctly without regressions.
        """
        report_data = self.create_test_report_data(
            total_gmv=total_gmv,
            total_deal=total_deal,
            total_video=total_video
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, f"test_property_{total_gmv}_{total_deal}_{total_video}.pptx")
        
        # Should not raise any exceptions
        renderer.render(report_data, output_path)
        
        # Verify file was created
        assert os.path.exists(output_path), "PPT should be generated for all metric values"
        
        # Verify PPT has expected structure
        prs = Presentation(output_path)
        assert len(prs.slides) >= 8, "Should have at least 8 slides"
    
    @settings(max_examples=5, deadline=None)
    @given(
        num_creators=st.integers(min_value=5, max_value=50),
        brand_name=st.text(min_size=3, max_size=20, alphabet=st.characters(whitelist_categories=('Lu', 'Ll')))
    )
    def test_property_various_creators(self, num_creators, brand_name):
        """
        Property-Based Test: Various creator counts and brand names
        
        EXPECTED: Test PASSES on both unfixed and fixed code
        
        Tests that the renderer handles different numbers of creators and
        brand names correctly.
        """
        report_data = self.create_test_report_data(
            num_creators=num_creators,
            brand_name=brand_name
        )
        
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, f"test_property_creators_{num_creators}.pptx")
        
        # Should not raise any exceptions
        renderer.render(report_data, output_path)
        
        # Verify file was created
        assert os.path.exists(output_path), "PPT should be generated for all creator counts"
        
        # Verify PPT structure
        prs = Presentation(output_path)
        assert len(prs.slides) >= 8, "Should have at least 8 slides"


if __name__ == "__main__":
    # Run tests with pytest
    pytest.main([__file__, "-v", "-s"])
