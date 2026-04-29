"""
PPT Renderer Section Images Bug Condition Exploration Test

This test demonstrates the bug where section_images are configured but not displayed
in the generated PowerPoint presentation because the image display code is commented out.

**CRITICAL**: This test MUST FAIL on unfixed code - failure confirms the bug exists.
**DO NOT attempt to fix the test or the code when it fails.**

The test will PASS after the fix is implemented, confirming the bug is resolved.
"""

import pytest
import os
import tempfile
from datetime import date
from pptx import Presentation
from PIL import Image

from app.services.ppt_renderer import PPTRenderer
from app.services.report_gen import ReportConfig, PerformanceMetrics, ReportData, CreatorRow
from app.services.brand_profile import BrandProfileData


class TestSectionImagesBugCondition:
    """
    Bug Condition Exploration Test - Section Images Not Displayed
    
    This test surfaces counterexamples that demonstrate the bug exists on UNFIXED code.
    When section_images are configured in ReportConfig, the generated PPT should display
    those images, but currently does NOT because the code is commented out.
    
    Expected behavior on UNFIXED code: Test FAILS (images missing)
    Expected behavior on FIXED code: Test PASSES (images present)
    """
    
    def setup_method(self):
        """Set up test fixtures and create test images."""
        # Create temporary directory for test images
        self.test_dir = tempfile.mkdtemp()
        self.test_images = {}
        
        # Create test images for each section
        section_keys = [
            "affiliate_performance_summary",
            "gmv_affiliate",
            "collaboration_metrics",
            "total_engagement"
        ]
        
        for section_key in section_keys:
            # Create a simple test image (100x100 blue square)
            img = Image.new('RGB', (100, 100), color='blue')
            img_path = os.path.join(self.test_dir, f"test_{section_key}.png")
            img.save(img_path)
            self.test_images[section_key] = img_path
        
        # Create uploads/images directory if it doesn't exist
        self.uploads_dir = os.path.join(os.getcwd(), "tiktok-affiliate-report", "uploads", "images")
        os.makedirs(self.uploads_dir, exist_ok=True)
        
        # Copy test images to uploads/images directory
        for section_key, img_path in self.test_images.items():
            dest_path = os.path.join(self.uploads_dir, f"test_{section_key}.png")
            img = Image.open(img_path)
            img.save(dest_path)
            self.test_images[section_key] = dest_path
    
    def teardown_method(self):
        """Clean up test images."""
        # Clean up test images from uploads/images
        for img_path in self.test_images.values():
            if os.path.exists(img_path):
                os.unlink(img_path)
        
        # Clean up temporary directory
        if os.path.exists(self.test_dir):
            import shutil
            shutil.rmtree(self.test_dir)
    
    def create_test_report_data(self, section_images=None):
        """Create test ReportData with section_images configured."""
        # Create test configuration
        config = ReportConfig(
            brand_name="Test Brand",
            batch_number="Test Batch 1",
            period_start=date(2024, 1, 1),
            period_end=date(2024, 1, 31),
            prev_gmv=10000000,
            section_images=section_images or {}
        )
        
        # Create test metrics
        metrics = PerformanceMetrics(
            total_approached=50,
            total_deal=30,
            total_posting=25,
            total_belum_posting=5,
            total_video=75,
            total_pesanan=100,
            total_produk_terjual=150,
            total_create_sale=20,
            total_gmv=15000000,
            hero_sku="Test SKU",
            gmv_change=5000000,
            gmv_change_pct=50.0
        )
        
        # Create test creator rows
        deal_rows = [
            CreatorRow(
                username=f"creator{i}",
                followers=10000 + (i * 1000),
                avg_gmv_month=500000 + (i * 100000),
                link_acc=f"https://tiktok.com/@creator{i}"
            )
            for i in range(10)
        ]
        
        # Create test brand profile
        brand_profile = BrandProfileData(
            name="Test Brand",
            sku_list=["SKU1", "SKU2", "SKU3"]
        )
        
        return ReportData(
            config=config,
            metrics=metrics,
            top_performers=deal_rows[:5],
            deal_rows=deal_rows,
            non_deal_rows=[],
            brand_profile=brand_profile
        )
    
    def count_images_in_slide(self, slide):
        """Count the number of image shapes in a slide."""
        image_count = 0
        for shape in slide.shapes:
            # Check if shape is a picture
            if hasattr(shape, 'image'):
                image_count += 1
        return image_count
    
    def test_performance_summary_image_missing(self):
        """
        Bug Condition: Performance Summary slide missing image
        
        EXPECTED ON UNFIXED CODE: Test FAILS (image missing)
        EXPECTED ON FIXED CODE: Test PASSES (image present)
        
        Counterexample: section_images["affiliate_performance_summary"] configured
        but image does NOT appear in generated PPT.
        """
        # Create report data with section_images configured
        section_images = {
            "affiliate_performance_summary": [self.test_images["affiliate_performance_summary"]]
        }
        report_data = self.create_test_report_data(section_images)
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_performance.pptx")
        renderer.render(report_data, output_path)
        
        # Extract Performance Summary slide (slide index 1, after cover)
        prs = Presentation(output_path)
        performance_slide = prs.slides[1]
        
        # Count images in slide
        image_count = self.count_images_in_slide(performance_slide)
        
        # ASSERTION: Image should be present (will FAIL on unfixed code)
        assert image_count > 0, (
            "COUNTEREXAMPLE FOUND: Performance Summary slide missing image despite "
            f"section_images['affiliate_performance_summary'] configured with path: "
            f"{self.test_images['affiliate_performance_summary']}. "
            "This confirms the bug exists - image display code is commented out."
        )
    
    def test_top10_image_missing(self):
        """
        Bug Condition: Top 10 Performer slide missing image
        
        EXPECTED ON UNFIXED CODE: Test FAILS (image missing)
        EXPECTED ON FIXED CODE: Test PASSES (image present)
        
        Counterexample: section_images["gmv_affiliate"] configured
        but image does NOT appear in generated PPT.
        """
        # Create report data with section_images configured
        section_images = {
            "gmv_affiliate": [self.test_images["gmv_affiliate"]]
        }
        report_data = self.create_test_report_data(section_images)
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_top10.pptx")
        renderer.render(report_data, output_path)
        
        # Extract Top 10 slide (slide index 3, after cover, performance, gmv_highlight)
        prs = Presentation(output_path)
        top10_slide = prs.slides[3]
        
        # Count images in slide
        image_count = self.count_images_in_slide(top10_slide)
        
        # ASSERTION: Image should be present (will FAIL on unfixed code)
        assert image_count > 0, (
            "COUNTEREXAMPLE FOUND: Top 10 Performer slide missing image despite "
            f"section_images['gmv_affiliate'] configured with path: "
            f"{self.test_images['gmv_affiliate']}. "
            "This confirms the bug exists - image display code is commented out."
        )
    
    def test_collaboration_image_missing(self):
        """
        Bug Condition: Collaboration Metrics slide missing image
        
        EXPECTED ON UNFIXED CODE: Test FAILS (image missing)
        EXPECTED ON FIXED CODE: Test PASSES (image present)
        
        Counterexample: section_images["collaboration_metrics"] configured
        but image does NOT appear in generated PPT.
        """
        # Create report data with section_images configured
        section_images = {
            "collaboration_metrics": [self.test_images["collaboration_metrics"]]
        }
        report_data = self.create_test_report_data(section_images)
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_collaboration.pptx")
        renderer.render(report_data, output_path)
        
        # Extract Collaboration slide (slide index 4, after cover, performance, gmv_highlight, top10)
        prs = Presentation(output_path)
        collaboration_slide = prs.slides[4]
        
        # Count images in slide
        image_count = self.count_images_in_slide(collaboration_slide)
        
        # ASSERTION: Image should be present (will FAIL on unfixed code)
        assert image_count > 0, (
            "COUNTEREXAMPLE FOUND: Collaboration Metrics slide missing image despite "
            f"section_images['collaboration_metrics'] configured with path: "
            f"{self.test_images['collaboration_metrics']}. "
            "This confirms the bug exists - image display code is commented out."
        )
    
    def test_engagement_image_missing(self):
        """
        Bug Condition: Engagement slide missing image
        
        EXPECTED ON UNFIXED CODE: Test FAILS (image missing)
        EXPECTED ON FIXED CODE: Test PASSES (image present)
        
        Counterexample: section_images["total_engagement"] configured
        but image does NOT appear in generated PPT.
        """
        # Create report data with section_images configured
        section_images = {
            "total_engagement": [self.test_images["total_engagement"]]
        }
        report_data = self.create_test_report_data(section_images)
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_engagement.pptx")
        renderer.render(report_data, output_path)
        
        # Extract Engagement slide (slide index 5, after cover, performance, gmv_highlight, top10, collaboration)
        prs = Presentation(output_path)
        engagement_slide = prs.slides[5]
        
        # Count images in slide
        image_count = self.count_images_in_slide(engagement_slide)
        
        # ASSERTION: Image should be present (will FAIL on unfixed code)
        assert image_count > 0, (
            "COUNTEREXAMPLE FOUND: Engagement slide missing image despite "
            f"section_images['total_engagement'] configured with path: "
            f"{self.test_images['total_engagement']}. "
            "This confirms the bug exists - image display code is commented out."
        )
    
    def test_all_section_images_missing(self):
        """
        Bug Condition: All four slides missing images
        
        EXPECTED ON UNFIXED CODE: Test FAILS (all images missing)
        EXPECTED ON FIXED CODE: Test PASSES (all images present)
        
        Comprehensive counterexample demonstrating the bug across all affected slides.
        """
        # Create report data with ALL section_images configured
        section_images = {
            "affiliate_performance_summary": [self.test_images["affiliate_performance_summary"]],
            "gmv_affiliate": [self.test_images["gmv_affiliate"]],
            "collaboration_metrics": [self.test_images["collaboration_metrics"]],
            "total_engagement": [self.test_images["total_engagement"]]
        }
        report_data = self.create_test_report_data(section_images)
        
        # Generate PPT
        renderer = PPTRenderer()
        output_path = os.path.join(self.test_dir, "test_all_sections.pptx")
        renderer.render(report_data, output_path)
        
        # Extract slides
        prs = Presentation(output_path)
        performance_slide = prs.slides[1]
        top10_slide = prs.slides[3]
        collaboration_slide = prs.slides[4]
        engagement_slide = prs.slides[5]
        
        # Count images in each slide
        performance_images = self.count_images_in_slide(performance_slide)
        top10_images = self.count_images_in_slide(top10_slide)
        collaboration_images = self.count_images_in_slide(collaboration_slide)
        engagement_images = self.count_images_in_slide(engagement_slide)
        
        # Collect missing images
        missing_slides = []
        if performance_images == 0:
            missing_slides.append("Performance Summary")
        if top10_images == 0:
            missing_slides.append("Top 10 Performer")
        if collaboration_images == 0:
            missing_slides.append("Collaboration Metrics")
        if engagement_images == 0:
            missing_slides.append("Engagement")
        
        # ASSERTION: All images should be present (will FAIL on unfixed code)
        assert len(missing_slides) == 0, (
            f"COUNTEREXAMPLES FOUND: {len(missing_slides)} slides missing images despite "
            f"section_images configured for all sections. Missing slides: {', '.join(missing_slides)}. "
            f"This confirms the bug exists - image display code is commented out in multiple slide methods."
        )


if __name__ == "__main__":
    # Run tests with pytest
    pytest.main([__file__, "-v", "-s"])
