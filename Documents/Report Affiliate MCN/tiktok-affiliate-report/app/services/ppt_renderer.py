"""
PPT Renderer — clean, simple slides using python-pptx.
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.services.report_gen import ReportData, format_rupiah, format_number

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x2C, 0x3E, 0x50)  # dark blue-gray
BLUE    = RGBColor(0x34, 0x98, 0xDB)  # bright blue
CORAL   = RGBColor(0xE7, 0x4C, 0x3C)  # coral red
LIGHT   = RGBColor(0xEC, 0xF0, 0xF1)  # light gray
LIGHT2  = RGBColor(0xD5, 0xDB, 0xDB)  # medium gray
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x2C, 0x3E, 0x50)
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
BORDER  = RGBColor(0xBD, 0xC3, 0xC7)
GREEN   = RGBColor(0x27, 0xAE, 0x60)  # positive green
RED     = RGBColor(0xE7, 0x4C, 0x3C)  # negative red

SW = 10.0   # slide width inches
SH = 7.5    # slide height inches
HEADER_H = 0.85


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _bg(slide, color=WHITE):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SW), Inches(SH))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _rect(slide, l, t, w, h, color, line_color=None, line_width=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def _txt(slide, l, t, w, h, text, size=12, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return txb


def _header(slide, title: str, subtitle: str = ""):
    _rect(slide, 0, 0, SW, HEADER_H, BLUE)
    _txt(slide, 0.3, 0.12, SW - 0.6, 0.5, title, size=24, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, 0.3, 0.58, SW - 0.6, 0.3, subtitle, size=11, color=LIGHT)


def _table(slide, headers: list, rows: list, left: float, top: float,
           col_widths: list, font_size: int = 10, max_height: float = None):
    """
    Render a table. Row height is calculated to fill available space.
    col_widths: list of floats in inches.
    max_height: max total table height in inches (auto-calculated if None).
    """
    if not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    # Calculate row height - use fixed reasonable height instead of filling space
    row_h = 0.28  # Fixed row height for consistent appearance

    total_w = Inches(sum(col_widths))
    total_h = Inches(row_h * n_rows)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        total_w, total_h
    )
    tbl = tbl_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(h)
        r.font.size = Pt(font_size)
        r.font.bold = True
        r.font.color.rgb = WHITE

    # Data rows
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            # First col center, second col left, rest center
            p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.color.rgb = BLACK

    return tbl_shape


def _metric_boxes(slide, items: list, top: float = 1.1, box_h: float = 1.2):
    """Render a row of metric boxes. items = [(label, value), ...]"""
    n = len(items)
    gap = 0.15
    total_gap = gap * (n - 1)
    box_w = (SW - 0.4 - total_gap) / n

    for i, (label, value) in enumerate(items):
        x = 0.2 + i * (box_w + gap)
        # Background box
        _rect(slide, x, top, box_w, box_h, WHITE, line_color=BORDER, line_width=1.5)
        # Top accent bar
        _rect(slide, x, top, box_w, 0.08, BLUE)
        
        # Value text - centered in upper portion
        value_top = top + 0.15
        value_height = box_h * 0.5
        _txt(slide, x, value_top, box_w, value_height, str(value),
             size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        
        # Label text - centered in lower portion
        label_top = top + box_h * 0.6
        label_height = box_h * 0.35
        _txt(slide, x, label_top, box_w, label_height, str(label),
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def _text_area(slide, l, t, w, h, text, size=12, color=BLACK, italic=False):
    """Multi-line text area."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    lines = str(text).split('\n')
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.italic = italic


def _add_image(slide, image_path: str, left: float, top: float, max_width: float = None, max_height: float = None):
    """
    Add image to slide with automatic proportional sizing.
    Returns True if successful.
    """
    if not os.path.isfile(image_path):
        return False
    
    try:
        from PIL import Image as PILImage
        
        # Get original image dimensions
        with PILImage.open(image_path) as pil_img:
            orig_w_px, orig_h_px = pil_img.size
        
        # Convert pixels to inches (assume 96 DPI)
        orig_w_in = orig_w_px / 96
        orig_h_in = orig_h_px / 96
        
        # Calculate scale to fit within max dimensions while preserving aspect ratio
        if max_width and max_height:
            scale_w = max_width / orig_w_in
            scale_h = max_height / orig_h_in
            scale = min(scale_w, scale_h, 1.5)  # Max 150% upscale
        elif max_width:
            scale = min(max_width / orig_w_in, 1.5)
        elif max_height:
            scale = min(max_height / orig_h_in, 1.5)
        else:
            scale = 1.0
        
        new_w = orig_w_in * scale
        new_h = orig_h_in * scale
        
        # Add image with calculated dimensions
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                width=Inches(new_w), height=Inches(new_h))
        return True
        
    except ImportError:
        # PIL not available - use python-pptx default sizing
        try:
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))
            # Try to scale if max dimensions provided
            if max_width and pic.width > Inches(max_width):
                ratio = Inches(max_width) / pic.width
                pic.width = Inches(max_width)
                pic.height = int(pic.height * ratio)
            if max_height and pic.height > Inches(max_height):
                ratio = Inches(max_height) / pic.height
                pic.height = Inches(max_height)
                pic.width = int(pic.width * ratio)
            return True
        except Exception:
            return False
            
    except Exception:
        # Fallback: try to add with default size
        try:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top))
            return True
        except Exception:
            return False


def _resolve_section_images(report_data: ReportData, section_key: str) -> list:
    """Resolve image paths for a section. Returns list of valid file paths."""
    # FIX 2: Helper to resolve section image paths
    images = (report_data.config.section_images or {}).get(section_key, [])
    if not images:
        return []
    
    resolved = []
    for img_id in images:
        # Try multiple path candidates
        for candidate in [
            img_id,
            os.path.join("uploads/images", img_id),
            os.path.join(os.getcwd(), "uploads/images", img_id),
        ]:
            if os.path.isfile(candidate):
                resolved.append(candidate)
                break
    return resolved


# ── PPTRenderer ───────────────────────────────────────────────────────────────

class PPTRenderer:

    def render(self, report_data: ReportData, output_path: str) -> str:
        prs = Presentation()
        prs.slide_width = Inches(SW)
        prs.slide_height = Inches(SH)
        blank = prs.slide_layouts[6]

        self._cover(prs, blank, report_data)
        self._performance(prs, blank, report_data)
        self._gmv_highlight(prs, blank, report_data)
        self._top10(prs, blank, report_data)
        self._collaboration(prs, blank, report_data)
        self._engagement(prs, blank, report_data)
        self._gmv_detail(prs, blank, report_data)
        self._insight_nextplan(prs, blank, report_data)

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        return output_path

    def render_multi_brand(self, brand_reports: dict, consolidated_config: dict, output_path: str) -> str:
        """
        Render consolidated multi-brand PowerPoint presentation.
        
        Args:
            brand_reports: Dict[brand_name, ReportData] - Individual brand report data
            consolidated_config: Dict with consolidated report configuration
            output_path: Output file path
        """
        prs = Presentation()
        prs.slide_width = Inches(SW)
        prs.slide_height = Inches(SH)
        blank = prs.slide_layouts[6]

        # Multi-brand slides
        self._multi_brand_cover(prs, blank, brand_reports, consolidated_config)
        self._multi_brand_executive_summary(prs, blank, brand_reports, consolidated_config)
        self._brand_comparison_slide(prs, blank, brand_reports, consolidated_config)
        
        # Individual brand sections
        for brand_name, report_data in brand_reports.items():
            self._brand_section_cover(prs, blank, brand_name, report_data)
            self._brand_performance_slide(prs, blank, brand_name, report_data)
            self._brand_gmv_highlight_slide(prs, blank, brand_name, report_data)
            self._brand_top_performers_slide(prs, blank, brand_name, report_data)
        
        # Consolidated insights
        self._multi_brand_insights_slide(prs, blank, brand_reports, consolidated_config)

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        return output_path

    # ── 1. Cover ──────────────────────────────────────────────────────────────

    def _cover(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide, NAVY)

        cfg = rd.config
        try:
            ps = cfg.period_start.strftime("%d %B %Y")
            pe = cfg.period_end.strftime("%d %B %Y")
        except Exception:
            ps, pe = str(cfg.period_start), str(cfg.period_end)

        # Accent line lebih tebal dan panjang
        _rect(slide, 3.0, 2.95, 4.0, 0.08, BLUE)

        _txt(slide, 0.5, 1.2, SW - 1, 0.5,
             "LAPORAN AFFILIATE TIKTOK", size=15, color=BLUE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 1.8, SW - 1, 1.3,
             cfg.brand_name.upper(), size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 3.2, SW - 1, 0.5,
             f"Batch {cfg.batch_number}", size=20, color=BLUE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 3.8, SW - 1, 0.4,
             f"Periode: {ps} — {pe}", size=13,
             color=LIGHT, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 6.9, SW - 1, 0.3,
             "Dokumen ini bersifat rahasia dan hanya untuk keperluan internal.",
             size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # ── 2. Performance Summary ────────────────────────────────────────────────

    def _performance(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        cfg = rd.config
        met = rd.metrics

        try:
            period_str = f"{cfg.period_start.strftime('%d %b')} s/d {cfg.period_end.strftime('%d %b %Y')}"
        except Exception:
            period_str = f"{cfg.period_start} s/d {cfg.period_end}"

        _header(slide, "Performance Summary", f"Batch {cfg.batch_number} — {period_str}")

        # Baris 1: 4 kolom - metrics dasar
        _metric_boxes(slide, [
            ("Total Deal", str(met.total_deal)),
            ("Total Posting", str(met.total_posting)),
            ("Belum Posting", str(met.total_belum_posting)),
            ("Total VT", format_number(met.total_video)),
        ], top=1.0, box_h=1.0)

        # Baris 2: 1 kolom - GMV (lebih besar dan menonjol)
        _metric_boxes(slide, [
            ("Total GMV", format_rupiah(met.total_gmv))
        ], top=2.15, box_h=1.2)

        # Baris 3: 3 kolom - metrics tambahan (hanya jika ada data)
        row3 = []
        if met.total_create_sale:
            row3.append(("Create Sale", str(met.total_create_sale)))
        if met.total_pesanan:
            row3.append(("Total Pesanan", format_number(met.total_pesanan)))
        if met.total_produk_terjual:
            row3.append(("Produk Terjual", format_number(met.total_produk_terjual)))
        
        if row3:
            _metric_boxes(slide, row3, top=3.5, box_h=1.0)
            table_top = 4.65
        else:
            table_top = 3.5

        # Summary table — fills remaining space
        rows = [
            ["1", "Total Akun Deal", f"{met.total_deal} akun"],
            ["2", "Total Akun Up VT", f"{met.total_posting} akun"],
            ["3", "Total Akun Belum Up VT", f"{met.total_belum_posting} akun"],
            ["4", "Total VT", format_number(met.total_video)],
            ["5", "GMV Batch Sebelumnya", format_rupiah(cfg.prev_gmv) if cfg.prev_gmv else "—"],
            ["6", f"GMV Batch {cfg.batch_number}", format_rupiah(met.total_gmv)],
        ]
        if met.total_create_sale:
            rows.append(["7", "Total Affiliate Create Sale", f"{met.total_create_sale} akun"])
        if met.total_produk_terjual or met.total_pesanan:
            rows.append(["8", "Total Produk Terjual / Pesanan",
                         f"{format_number(met.total_produk_terjual)} / {format_number(met.total_pesanan)}"])

        _table(slide, ["No", "Keterangan", "Data"],
               rows, left=0.3, top=table_top,
               col_widths=[0.5, 5.7, 3.3],
               font_size=10, max_height=SH - table_top - 0.15)

        # Section images - positioned below table to avoid overlap
        images = _resolve_section_images(rd, "affiliate_performance_summary")
        if images:
            # Position image below table with small gap
            # Table ends at approximately table_top + (num_rows * 0.28) + 0.15
            # For safety, position at 6.5 to ensure no overlap
            for img_path in images[:1]:
                _add_image(slide, img_path, 0.3, 6.5, max_width=3.0, max_height=0.9)

    # ── 3. GMV Highlight ─────────────────────────────────────────────────────

    def _gmv_highlight(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "GMV Highlight")

        m = rd.metrics
        cfg = rd.config

        # GMV box dengan warna coral untuk highlight
        _rect(slide, 1.5, 1.1, 7.0, 2.3, CORAL)
        _txt(slide, 1.5, 1.25, 7.0, 0.4,
             f"Total GMV Batch {cfg.batch_number}", size=12, color=LIGHT, align=PP_ALIGN.CENTER)
        _txt(slide, 1.5, 1.65, 7.0, 1.0,
             format_rupiah(m.total_gmv), size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        if m.gmv_change is not None:
            sign = "▲" if m.gmv_change >= 0 else "▼"
            c = GREEN if m.gmv_change >= 0 else WHITE
            direction = "naik" if m.gmv_change >= 0 else "turun"
            _txt(slide, 1.5, 2.7, 7.0, 0.45,
                 f"{sign} {format_rupiah(abs(m.gmv_change))} ({abs(m.gmv_change_pct):.1f}% {direction} vs batch sebelumnya)",
                 size=12, color=c, align=PP_ALIGN.CENTER)

        if rd.brand_profile and rd.brand_profile.sku_list:
            _rect(slide, 1.5, 3.6, 7.0, 0.9, LIGHT, line_color=BORDER, line_width=1.5)
            _txt(slide, 1.5, 3.65, 7.0, 0.28, "PRODUK / SKU",
                 size=9, color=GRAY, align=PP_ALIGN.CENTER)
            _txt(slide, 1.5, 3.93, 7.0, 0.5,
                 "  •  ".join(rd.brand_profile.sku_list[:5]),
                 size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

        # Add section images if available (ada space di bawah)
        images = _resolve_section_images(rd, "gmv_batch")
        if images:
            img_top = 4.7
            max_img_width = 7.0
            max_img_height = 2.5
            for img_path in images[:2]:  # Max 2 images
                _add_image(slide, img_path, 1.5, img_top, max_width=max_img_width, max_height=max_img_height)
                img_top += 0.1  # Small offset if multiple images

    # ── 4. Top 10 ─────────────────────────────────────────────────────────────

    def _top10(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Top 10 Performer", "Berdasarkan GMV tertinggi")

        top = rd.top_performers[:10]
        if not top:
            _txt(slide, 0.5, 2, 9, 0.5, "Tidak ada data.", size=12, color=GRAY)
            return

        rows = [[str(i + 1), r.username or "—", format_rupiah(r.avg_gmv_month or 0)]
                for i, r in enumerate(top)]

        # Tabel lebih besar dan centered
        _table(slide, ["#", "Username", "GMV"],
               rows, left=1.0, top=1.0,
               col_widths=[0.6, 5.5, 2.9],
               font_size=11, max_height=SH - 1.0 - 0.15)

        # Section images - positioned below table to avoid overlap
        images = _resolve_section_images(rd, "gmv_affiliate")
        if images:
            # Position image below table (table ends around 6.5-7.0)
            # Place at bottom of slide with small margin
            for img_path in images[:1]:
                _add_image(slide, img_path, 1.0, 6.8, max_width=8.0, max_height=0.6)

    # ── 5. Collaboration ──────────────────────────────────────────────────────

    def _collaboration(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Collaboration Metrics")

        m = rd.metrics
        rasio = f"{m.total_deal / m.total_approached * 100:.1f}%" if m.total_approached > 0 else "0%"

        # Metric boxes centered dan lebih besar
        _metric_boxes(slide, [
            ("Total Approached", format_number(m.total_approached)),
            ("Total Deal", format_number(m.total_deal)),
            ("Rasio Konversi", rasio),
        ], top=2.0, box_h=2.5)

        # Section images - positioned below metric boxes (similar to GMV Highlight pattern)
        images = _resolve_section_images(rd, "collaboration_metrics")
        if images:
            # Metric boxes end at top=2.0 + box_h=2.5 = 4.5
            # Position images below with small gap, similar to GMV Highlight slide
            img_top = 4.7
            max_img_width = 7.0
            max_img_height = 2.5
            for img_path in images[:2]:
                _add_image(slide, img_path, 1.5, img_top, max_width=max_img_width, max_height=max_img_height)
                img_top += 0.1  # Small offset if multiple images

    # ── 6. Engagement ─────────────────────────────────────────────────────────

    def _engagement(self, prs, layout, rd: ReportData):
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Total Engagement")

        eng = rd.config.engagement
        if eng:
            _metric_boxes(slide, [
                ("Total Views",    f"{eng.total_views:,}".replace(",", ".")),
                ("Total Likes",    f"{eng.total_likes:,}".replace(",", ".")),
                ("Total Comments", f"{eng.total_comments:,}".replace(",", ".")),
            ], top=1.0, box_h=1.2)

            # Top 5 creator dengan breakdown per video
            if eng.top_creators:
                _txt(slide, 0.4, 2.4, SW - 0.8, 0.35,
                     "Top 5 Creator Engagement", size=12, bold=True, color=BLUE)
                
                # DEBUG: Print creator info
                print(f"[PPT DEBUG] Processing {len(eng.top_creators)} creators")
                for c in eng.top_creators[:5]:
                    print(f"[PPT DEBUG]   - {c.username}: {len(c.videos) if c.videos else 0} videos")
                
                # ALWAYS use 8-column format if we have engagement data
                # This ensures consistent layout regardless of video count
                headers = ["#", "Username", "Views (vid1)", "Likes (vid1)", "Comments (vid1)", 
                          "Views (vid2)", "Likes (vid2)", "Comments (vid2)"]
                table_rows = []
                
                for i, c in enumerate(eng.top_creators[:5]):
                    if c.videos and len(c.videos) >= 2:
                        # Creator dengan 2+ video
                        vid1 = c.videos[0]
                        vid2 = c.videos[1]
                        table_rows.append([
                            str(i + 1),
                            f"@{c.username}",
                            f"{vid1.views:,}".replace(",", "."),
                            f"{vid1.likes:,}".replace(",", "."),
                            f"{vid1.comments:,}".replace(",", "."),
                            f"{vid2.views:,}".replace(",", "."),
                            f"{vid2.likes:,}".replace(",", "."),
                            f"{vid2.comments:,}".replace(",", "."),
                        ])
                    elif c.videos and len(c.videos) == 1:
                        # Creator dengan 1 video - kolom vid2 kosong
                        vid1 = c.videos[0]
                        table_rows.append([
                            str(i + 1),
                            f"@{c.username}",
                            f"{vid1.views:,}".replace(",", "."),
                            f"{vid1.likes:,}".replace(",", "."),
                            f"{vid1.comments:,}".replace(",", "."),
                            "—", "—", "—",
                        ])
                    else:
                        # Creator tanpa breakdown video - gunakan total
                        table_rows.append([
                            str(i + 1),
                            f"@{c.username}",
                            f"{c.total_views:,}".replace(",", "."),
                            f"{c.total_likes:,}".replace(",", "."),
                            f"{c.total_comments:,}".replace(",", "."),
                            "—", "—", "—",
                        ])
                
                _table(slide, headers, table_rows, left=0.3, top=2.85,
                       col_widths=[0.3, 1.8, 1.2, 1.0, 1.0, 1.2, 1.0, 1.0],
                       font_size=8, max_height=SH - 2.85 - 0.2)
        else:
            _txt(slide, 0.5, 2.5, SW - 1, 0.5,
                 "Data engagement tidak tersedia.", size=13, color=GRAY, align=PP_ALIGN.CENTER)

        # Section images - positioned below engagement content
        images = _resolve_section_images(rd, "total_engagement")
        if images:
            # Position depends on whether top_creators table exists
            # If table exists: position below table (around 6.5-7.0)
            # If no table: more space available (around 3.5)
            img_top = 6.5 if eng and eng.top_creators else 3.5
            max_img_height = 0.8 if eng and eng.top_creators else 3.5
            for img_path in images[:1]:
                _add_image(slide, img_path, 0.4, img_top, max_width=4.0, max_height=max_img_height)

    # ── 7. GMV Detail ─────────────────────────────────────────────────────────

    def _gmv_detail(self, prs, layout, rd: ReportData):
        rows_with_gmv = sorted(
            [r for r in rd.deal_rows if r.avg_gmv_month and r.avg_gmv_month > 0],
            key=lambda r: r.avg_gmv_month,
            reverse=True,
        )

        if not rows_with_gmv:
            slide = prs.slides.add_slide(layout)
            _bg(slide)
            _header(slide, f"GMV Affiliate Batch {rd.config.batch_number}")
            _txt(slide, 0.5, 2, SW - 1, 0.5, "Tidak ada data GMV.", size=12, color=GRAY)
            return

        # Layout: 2 columns side by side, each with its own table
        # Each slide shows up to 20 rows (10 per column)
        ROWS_PER_COL = 10
        ROWS_PER_SLIDE = ROWS_PER_COL * 2  # 20 per slide

        # Column layout
        COL_L = 0.2          # left column x
        COL_R = 5.2          # right column x
        COL_W = [0.5, 2.8, 1.5]  # #, Username, GMV — total = 4.8
        TABLE_TOP = 1.0

        for slide_idx, start in enumerate(range(0, len(rows_with_gmv), ROWS_PER_SLIDE)):
            chunk = rows_with_gmv[start:start + ROWS_PER_SLIDE]
            slide = prs.slides.add_slide(layout)
            _bg(slide)

            subtitle = f"Total: {format_rupiah(rd.metrics.total_gmv)}"
            if slide_idx > 0:
                subtitle += f" (lanjutan {slide_idx + 1})"
            _header(slide, f"GMV Affiliate Batch {rd.config.batch_number}", subtitle)

            left_rows = chunk[:ROWS_PER_COL]
            right_rows = chunk[ROWS_PER_COL:]

            # Left column table
            if left_rows:
                data_l = [
                    [str(start + i + 1), r.username or "—", format_rupiah(r.avg_gmv_month or 0)]
                    for i, r in enumerate(left_rows)
                ]
                _table(slide, ["#", "Username", "GMV"],
                       data_l, left=COL_L, top=TABLE_TOP,
                       col_widths=COL_W, font_size=9)

            # Right column table
            if right_rows:
                data_r = [
                    [str(start + ROWS_PER_COL + i + 1), r.username or "—", format_rupiah(r.avg_gmv_month or 0)]
                    for i, r in enumerate(right_rows)
                ]
                _table(slide, ["#", "Username", "GMV"],
                       data_r, left=COL_R, top=TABLE_TOP,
                       col_widths=COL_W, font_size=9)

            # Bottom info bar (only first slide)
            if slide_idx == 0:
                info_parts = [f"Total GMV: {format_rupiah(rd.metrics.total_gmv)}"]
                if rd.metrics.total_create_sale:
                    info_parts.append(f"Create Sale: {rd.metrics.total_create_sale} akun")
                _rect(slide, 0.2, SH - 0.6, SW - 0.4, 0.45, LIGHT, line_color=BORDER, line_width=1.5)
                _txt(slide, 0.3, SH - 0.55, SW - 0.6, 0.35,
                     "  •  ".join(info_parts),
                     size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # ── 8. Insight & Next Plan ────────────────────────────────────────────────

    def _insight_nextplan(self, prs, layout, rd: ReportData):
        cfg = rd.config

        for title, text_val in [("Insight", cfg.insight), ("Next Plan", cfg.next_plan)]:
            slide = prs.slides.add_slide(layout)
            _bg(slide)
            _header(slide, title)

            text = text_val or "Belum diisi."
            is_empty = not bool(text_val)

            # Background box untuk text area
            _rect(slide, 0.4, 1.1, SW - 0.8, SH - 1.4, LIGHT, line_color=BLUE, line_width=2)
            
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(SW - 1.0), Inches(SH - 1.5))
            tf = txb.text_frame
            tf.word_wrap = True
            
            lines = str(text).split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                    
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                
                # Add bullet point
                if not is_empty:
                    p.text = f"• {line}"
                else:
                    p.text = line
                    
                # Format text
                for run in p.runs:
                    run.font.size = Pt(15)  # Larger font size
                    run.font.color.rgb = GRAY if is_empty else BLACK
                    run.font.italic = is_empty

    # ── Multi-Brand Slide Methods ─────────────────────────────────────────────

    def _multi_brand_cover(self, prs, layout, brand_reports: dict, consolidated_config: dict):
        """Multi-brand cover slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide, NAVY)

        batch = consolidated_config.get('batch_number', 'Batch 1')
        brand_names = list(brand_reports.keys())
        brand_list = " • ".join(brand_names[:3])  # Show first 3 brands
        if len(brand_names) > 3:
            brand_list += f" • +{len(brand_names) - 3} more"

        try:
            ps = consolidated_config.get('period_start', '').strftime("%d %B %Y") if consolidated_config.get('period_start') else ''
            pe = consolidated_config.get('period_end', '').strftime("%d %B %Y") if consolidated_config.get('period_end') else ''
            period_str = f"{ps} — {pe}" if ps and pe else ""
        except Exception:
            period_str = ""

        # Accent line
        _rect(slide, 2.5, 2.95, 5.0, 0.08, BLUE)

        _txt(slide, 0.5, 1.0, SW - 1, 0.5,
             "LAPORAN AFFILIATE TIKTOK", size=15, color=BLUE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 1.6, SW - 1, 1.0,
             "MULTI-BRAND REPORT", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 2.6, SW - 1, 0.6,
             brand_list, size=16, color=BLUE, align=PP_ALIGN.CENTER)
        _txt(slide, 0.5, 3.3, SW - 1, 0.5,
             f"Batch {batch}", size=18, color=BLUE, align=PP_ALIGN.CENTER)
        if period_str:
            _txt(slide, 0.5, 3.9, SW - 1, 0.4,
                 f"Periode: {period_str}", size=13, color=LIGHT, align=PP_ALIGN.CENTER)
        
        # Summary stats
        total_creators = sum(len(rd.deal_rows) for rd in brand_reports.values())
        _txt(slide, 0.5, 4.5, SW - 1, 0.4,
             f"{len(brand_names)} brands • {total_creators} creators", size=12, color=LIGHT, align=PP_ALIGN.CENTER)
        
        _txt(slide, 0.5, 6.9, SW - 1, 0.3,
             "Dokumen ini bersifat rahasia dan hanya untuk keperluan internal.",
             size=9, color=GRAY, align=PP_ALIGN.CENTER)

    def _multi_brand_executive_summary(self, prs, layout, brand_reports: dict, consolidated_config: dict):
        """Executive summary slide with aggregated metrics."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Executive Summary", f"{len(brand_reports)} brands consolidated")

        # Aggregate metrics
        total_deal = sum(rd.metrics.total_deal for rd in brand_reports.values())
        total_posting = sum(rd.metrics.total_posting for rd in brand_reports.values())
        total_video = sum(rd.metrics.total_video for rd in brand_reports.values())
        total_gmv = sum(rd.metrics.total_gmv for rd in brand_reports.values())

        # Metrics boxes
        _metric_boxes(slide, [
            ("Total Brands", str(len(brand_reports))),
            ("Total Deal", str(total_deal)),
            ("Total Posting", str(total_posting)),
            ("Total VT", format_number(total_video)),
        ], top=1.0, box_h=1.0)

        _metric_boxes(slide, [
            ("Total GMV", format_rupiah(total_gmv))
        ], top=2.2, box_h=1.2)

        # Summary table
        rows = [
            ["1", "Jumlah Brand", f"{len(brand_reports)} brands"],
            ["2", "Total Akun Deal", f"{total_deal} akun"],
            ["3", "Total Akun Posting", f"{total_posting} akun"],
            ["4", "Total Video (VT)", format_number(total_video)],
            ["5", "Total GMV", format_rupiah(total_gmv)],
        ]

        _table(slide, ["No", "Metrik", "Total"],
               rows, left=0.5, top=3.6,
               col_widths=[0.5, 4.5, 4.0],
               font_size=11)

    def _brand_comparison_slide(self, prs, layout, brand_reports: dict, consolidated_config: dict):
        """Brand comparison slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Brand Comparison", "Performance comparison across brands")

        if len(brand_reports) <= 1:
            _txt(slide, 0.5, 2.5, SW - 1, 0.5,
                 "Perbandingan memerlukan minimal 2 brand.", size=13, color=GRAY, align=PP_ALIGN.CENTER)
            return

        # Comparison table
        rows = []
        for brand_name, report_data in brand_reports.items():
            m = report_data.metrics
            avg_gmv = (m.total_gmv / m.total_deal) if m.total_deal > 0 else 0
            rows.append([
                brand_name,
                str(m.total_deal),
                str(m.total_posting),
                format_rupiah(m.total_gmv),
                format_rupiah(avg_gmv)
            ])

        # Sort by GMV descending
        rows.sort(key=lambda x: float(x[3].replace('Rp', '').replace('.', '').replace(',', '')) if 'Rp' in x[3] else 0, reverse=True)

        _table(slide, ["Brand", "Deal", "Posting", "Total GMV", "Avg GMV/Creator"],
               rows, left=0.3, top=1.2,
               col_widths=[2.0, 1.5, 1.5, 2.5, 2.2],
               font_size=10)

    def _brand_section_cover(self, prs, layout, brand_name: str, report_data: ReportData):
        """Brand section cover slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide, LIGHT)

        # Brand identification header
        _rect(slide, 0, 0, SW, 1.5, BLUE)
        _txt(slide, 0.5, 0.3, SW - 1, 0.9,
             f"BRAND: {brand_name}", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Brand stats
        creators_count = len(report_data.deal_rows)
        total_gmv = report_data.metrics.total_gmv

        _metric_boxes(slide, [
            ("Creators", str(creators_count)),
            ("Total GMV", format_rupiah(total_gmv)),
        ], top=2.5, box_h=1.5)

        # Brand profile info if available
        if report_data.brand_profile and report_data.brand_profile.sku_list:
            _rect(slide, 1.0, 4.5, SW - 2.0, 1.2, WHITE, line_color=BORDER, line_width=1.5)
            _txt(slide, 1.0, 4.6, SW - 2.0, 0.3, "PRODUK / SKU",
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)
            _txt(slide, 1.0, 4.9, SW - 2.0, 0.7,
                 "  •  ".join(report_data.brand_profile.sku_list[:4]),
                 size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    def _brand_performance_slide(self, prs, layout, brand_name: str, report_data: ReportData):
        """Brand-specific performance slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        
        cfg = report_data.config
        met = report_data.metrics

        try:
            period_str = f"{cfg.period_start.strftime('%d %b')} s/d {cfg.period_end.strftime('%d %b %Y')}"
        except Exception:
            period_str = f"{cfg.period_start} s/d {cfg.period_end}"

        _header(slide, f"Performance - {brand_name}", f"Batch {cfg.batch_number} — {period_str}")

        # Brand identification box
        _rect(slide, 0.3, 0.9, SW - 0.6, 0.3, LIGHT, line_color=BLUE, line_width=1)
        _txt(slide, 0.3, 0.95, SW - 0.6, 0.2,
             f"Brand: {brand_name} • {len(report_data.deal_rows)} creators",
             size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

        # Metrics (same as regular performance slide but with brand context)
        _metric_boxes(slide, [
            ("Deal", str(met.total_deal)),
            ("Posting", str(met.total_posting)),
            ("Belum Posting", str(met.total_belum_posting)),
            ("Total VT", format_number(met.total_video)),
        ], top=1.4, box_h=1.0)

        _metric_boxes(slide, [
            ("GMV", format_rupiah(met.total_gmv))
        ], top=2.6, box_h=1.2)

    def _brand_gmv_highlight_slide(self, prs, layout, brand_name: str, report_data: ReportData):
        """Brand-specific GMV highlight slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, f"GMV Highlight - {brand_name}")

        m = report_data.metrics
        cfg = report_data.config

        # GMV box with brand identification
        _rect(slide, 1.5, 1.1, 7.0, 2.5, CORAL)
        _txt(slide, 1.5, 1.25, 7.0, 0.3,
             f"{brand_name}", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, 1.5, 1.55, 7.0, 0.4,
             f"GMV Batch {cfg.batch_number}", size=12, color=LIGHT, align=PP_ALIGN.CENTER)
        _txt(slide, 1.5, 1.95, 7.0, 1.0,
             format_rupiah(m.total_gmv), size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Brand profile SKUs
        if report_data.brand_profile and report_data.brand_profile.sku_list:
            _rect(slide, 1.5, 3.8, 7.0, 0.9, LIGHT, line_color=BORDER, line_width=1.5)
            _txt(slide, 1.5, 3.85, 7.0, 0.28, "PRODUK / SKU",
                 size=9, color=GRAY, align=PP_ALIGN.CENTER)
            _txt(slide, 1.5, 4.13, 7.0, 0.5,
                 "  •  ".join(report_data.brand_profile.sku_list[:5]),
                 size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    def _brand_top_performers_slide(self, prs, layout, brand_name: str, report_data: ReportData):
        """Brand-specific top performers slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, f"Top Performers - {brand_name}", "Top 10 berdasarkan GMV")

        top = report_data.top_performers[:10]
        if not top:
            _txt(slide, 0.5, 2, 9, 0.5, "Tidak ada data.", size=12, color=GRAY)
            return

        rows = [[str(i + 1), r.username or "—", format_rupiah(r.avg_gmv_month or 0)]
                for i, r in enumerate(top)]

        _table(slide, ["#", "Username", "GMV"],
               rows, left=1.0, top=1.0,
               col_widths=[0.6, 5.5, 2.9],
               font_size=11)

    def _multi_brand_insights_slide(self, prs, layout, brand_reports: dict, consolidated_config: dict):
        """Multi-brand insights and next plan slide."""
        slide = prs.slides.add_slide(layout)
        _bg(slide)
        _header(slide, "Multi-Brand Insights")

        # Generate automatic insights
        insights = []
        
        # Best performing brand
        best_brand = max(brand_reports.items(), key=lambda x: x[1].metrics.total_gmv)
        insights.append(f"Brand dengan performa terbaik: {best_brand[0]} dengan GMV {format_rupiah(best_brand[1].metrics.total_gmv)}")
        
        # Total summary
        total_gmv = sum(rd.metrics.total_gmv for rd in brand_reports.values())
        total_creators = sum(len(rd.deal_rows) for rd in brand_reports.values())
        insights.append(f"Total {len(brand_reports)} brand dengan {total_creators} creators menghasilkan GMV {format_rupiah(total_gmv)}")
        
        # Conversion rates
        conversion_rates = {}
        for brand_name, report_data in brand_reports.items():
            if report_data.metrics.total_approached > 0:
                conversion_rates[brand_name] = report_data.metrics.total_deal / report_data.metrics.total_approached
        
        if conversion_rates:
            best_conversion_brand = max(conversion_rates.items(), key=lambda x: x[1])
            insights.append(f"Brand dengan konversi terbaik: {best_conversion_brand[0]} ({best_conversion_brand[1]*100:.1f}%)")

        # Display insights
        _rect(slide, 0.4, 1.1, SW - 0.8, 2.5, LIGHT, line_color=BLUE, line_width=2)
        
        insight_text = "\n\n".join(insights)
        _text_area(slide, 0.5, 1.2, SW - 1.0, 2.3, insight_text, size=13, color=BLACK)

        # Next plan
        _rect(slide, 0.4, 4.0, SW - 0.8, 2.5, WHITE, line_color=CORAL, line_width=2)
        _txt(slide, 0.5, 4.1, SW - 1.0, 0.3, "NEXT PLAN", size=12, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
        
        next_plan = consolidated_config.get('next_plan', 'Fokus pada optimasi brand dengan performa terbaik dan perbaikan strategi untuk brand dengan konversi rendah.')
        _text_area(slide, 0.5, 4.4, SW - 1.0, 2.0, next_plan, size=12, color=BLACK)
